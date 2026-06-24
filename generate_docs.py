import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ==========================================
# CONSTANTS & DESIGN SYSTEM
# ==========================================
BG_DARK = RGBColor(11, 15, 25)       # #0b0f19
TEXT_WHITE = RGBColor(255, 255, 255) # #ffffff
TEXT_MUTED = RGBColor(156, 163, 175) # #9ca3af
COLOR_PRIMARY = RGBColor(99, 102, 241)  # #6366f1 (Indigo)
COLOR_SECONDARY = RGBColor(6, 182, 212) # #06b6d4 (Cyan)
COLOR_ACCENT = RGBColor(168, 85, 247)   # #a855f7 (Purple)
COLOR_CARD = RGBColor(17, 24, 39)       # #111827
COLOR_BORDER = RGBColor(31, 41, 55)     # #1f2937

# ==========================================
# POWERPOINT GENERATOR
# ==========================================
def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # completely blank layout

    # Helper to set slide background
    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_DARK

    # Helper to add a colored top bar
    def add_top_bar(slide):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_SECONDARY
        shape.line.fill.background()

    # Helper to add footer
    def add_footer(slide, slide_num):
        # Footer text
        txBox = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(10), Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "FOODFLOW 4.1  |  B.Tech AI & ML Internship Project"
        p.font.name = "Inter"
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_MUTED
        
        # Slide number
        numBox = slide.shapes.add_textbox(Inches(12.0), Inches(7.0), Inches(0.8), Inches(0.4))
        num_p = numBox.text_frame.paragraphs[0]
        num_p.text = str(slide_num)
        num_p.font.name = "Inter"
        num_p.font.size = Pt(9)
        num_p.font.color.rgb = TEXT_MUTED
        num_p.alignment = PP_ALIGN.RIGHT

    # Helper to add standard header
    def add_header(slide, title_text):
        txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Outfit"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
        # Add accent line under title
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.1), Inches(4), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_PRIMARY
        line.line.fill.background()

    # ------------------------------------------
    # SLIDE 1: Title Slide
    # ------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_background(slide1)
    
    # Large Decorative Glow box
    glow = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(2), Inches(9.33), Inches(3.5))
    glow.fill.solid()
    glow.fill.fore_color.rgb = COLOR_CARD
    glow.line.color.rgb = COLOR_PRIMARY
    glow.line.width = Pt(1.5)
    
    # Title badge
    badgeBox = slide1.shapes.add_textbox(Inches(2.5), Inches(2.2), Inches(8.33), Inches(0.4))
    badge_p = badgeBox.text_frame.paragraphs[0]
    badge_p.text = "INTERNSHIP PROJECT PRESENTATION"
    badge_p.font.name = "Inter"
    badge_p.font.size = Pt(10)
    badge_p.font.bold = True
    badge_p.font.color.rgb = COLOR_SECONDARY
    badge_p.alignment = PP_ALIGN.CENTER

    # Title text
    titleBox = slide1.shapes.add_textbox(Inches(2.5), Inches(2.7), Inches(8.33), Inches(1.0))
    title_p = titleBox.text_frame.paragraphs[0]
    title_p.text = "FOODFLOW 4.1"
    title_p.font.name = "Outfit"
    title_p.font.size = Pt(48)
    title_p.font.bold = True
    title_p.font.color.rgb = TEXT_WHITE
    title_p.alignment = PP_ALIGN.CENTER
    
    # Subtitle text
    subBox = slide1.shapes.add_textbox(Inches(2.5), Inches(3.8), Inches(8.33), Inches(0.8))
    sub_p = subBox.text_frame.paragraphs[0]
    sub_p.text = "Cloud-Based Multi-Vendor Food Ordering SaaS Ecosystem"
    sub_p.font.name = "Inter"
    sub_p.font.size = Pt(16)
    sub_p.font.color.rgb = TEXT_MUTED
    sub_p.alignment = PP_ALIGN.CENTER

    # Presenter Information Left
    infoBoxL = slide1.shapes.add_textbox(Inches(2), Inches(5.8), Inches(4.5), Inches(1.2))
    tfL = infoBoxL.text_frame
    tfL.word_wrap = True
    pL1 = tfL.paragraphs[0]
    pL1.text = "SUBMITTED BY"
    pL1.font.name = "Inter"
    pL1.font.size = Pt(9)
    pL1.font.bold = True
    pL1.font.color.rgb = TEXT_MUTED
    
    pL2 = tfL.add_paragraph()
    pL2.text = "Sreeharan M Anilkumar"
    pL2.font.name = "Outfit"
    pL2.font.size = Pt(14)
    pL2.font.bold = True
    pL2.font.color.rgb = TEXT_WHITE
    
    pL3 = tfL.add_paragraph()
    pL3.text = "Reg No: MEC23AIM006\nB.Tech, Artificial Intelligence & ML"
    pL3.font.name = "Inter"
    pL3.font.size = Pt(10)
    pL3.font.color.rgb = TEXT_MUTED

    # Presenter Information Right
    infoBoxR = slide1.shapes.add_textbox(Inches(6.83), Inches(5.8), Inches(4.5), Inches(1.2))
    tfR = infoBoxR.text_frame
    tfR.word_wrap = True
    pR1 = tfR.paragraphs[0]
    pR1.text = "HOST ORGANIZATION"
    pR1.font.name = "Inter"
    pR1.font.size = Pt(9)
    pR1.font.bold = True
    pR1.font.color.rgb = TEXT_MUTED
    
    pR2 = tfR.add_paragraph()
    pR2.text = "Bairuha Tech"
    pR2.font.name = "Outfit"
    pR2.font.size = Pt(14)
    pR2.font.bold = True
    pR2.font.color.rgb = TEXT_WHITE
    
    pR3 = tfR.add_paragraph()
    pR3.text = "Mentors: Ambika Raj, Ressaan\nPlatform: Next.js & NestJS"
    pR3.font.name = "Inter"
    pR3.font.size = Pt(10)
    pR3.font.color.rgb = TEXT_MUTED

    add_top_bar(slide1)
    add_footer(slide1, 1)

    # ------------------------------------------
    # SLIDE 2: Project Overview & Objectives
    # ------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_background(slide2)
    add_top_bar(slide2)
    add_header(slide2, "Project Overview & Objectives")
    add_footer(slide2, 2)

    # Left content box (Challenges)
    leftBox = slide2.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.8))
    tf = leftBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "The Industry Challenge"
    p.font.name = "Outfit"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_SECONDARY
    
    bullets = [
      ("Fragmented Workflows", "Local restaurant merchants struggle with isolated menus, complex ordering interfaces, and offline accounting inefficiencies."),
      ("Row-Level Data Security", "Lack of proper multi-tenant boundaries on sharing databases lead to cross-tenant data leakage risks."),
      ("Static Static Tracking", "Customers face static, non-realtime food dispatch experiences without dynamic downloadable receipts.")
    ]
    for b_title, b_desc in bullets:
        p = tf.add_paragraph()
        p.text = f"\n• {b_title}: "
        p.font.name = "Inter"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
        p2 = tf.add_paragraph()
        p2.text = b_desc
        p2.font.name = "Inter"
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = TEXT_MUTED

    # Right Content (Two customized cards)
    # Card 1
    card1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.8), Inches(2.2))
    card1.fill.solid()
    card1.fill.fore_color.rgb = COLOR_CARD
    card1.line.color.rgb = COLOR_PRIMARY
    card1.line.width = Pt(1)
    
    tb1 = slide2.shapes.add_textbox(Inches(7.0), Inches(1.9), Inches(5.4), Inches(2.0))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "Real-Time WebSocket Streams"
    p1.font.name = "Outfit"
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_PRIMARY
    
    p1_desc = tf1.add_paragraph()
    p1_desc.text = "\nEnables instantaneous state synchronization between customers and vendors using Socket.IO, eliminating manual browser reloads."
    p1_desc.font.name = "Inter"
    p1_desc.font.size = Pt(12)
    p1_desc.font.color.rgb = TEXT_MUTED

    # Card 2
    card2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(4.3), Inches(5.8), Inches(2.2))
    card2.fill.solid()
    card2.fill.fore_color.rgb = COLOR_CARD
    card2.line.color.rgb = COLOR_ACCENT
    card2.line.width = Pt(1)
    
    tb2 = slide2.shapes.add_textbox(Inches(7.0), Inches(4.4), Inches(5.4), Inches(2.0))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Decoupled Multi-Portal Strategy"
    p2.font.name = "Outfit"
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_ACCENT
    
    p2_desc = tf2.add_paragraph()
    p2_desc.text = "\nDedicated portals specifically tailored to the workflows of Customer, Restaurant Vendor Admins, and platform Super Administrators."
    p2_desc.font.name = "Inter"
    p2_desc.font.size = Pt(12)
    p2_desc.font.color.rgb = TEXT_MUTED

    # ------------------------------------------
    # SLIDE 3: Technology Stack
    # ------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_background(slide3)
    add_top_bar(slide3)
    add_header(slide3, "Enterprise Technology Stack")
    add_footer(slide3, 3)

    # Grid of 4 categories
    tech_data = [
      ("Frontend Application", "Next.js 16/15 (App Router)\nTypeScript & React 19\nTailwind CSS v4\nFramer Motion (Animations)\nRecharts (Analytics Charts)", COLOR_SECONDARY, Inches(0.6), Inches(1.8)),
      ("Backend Server API", "NestJS 11 Gateway\nNode.js Runtime\nSocket.IO (WebSockets)\nRxJS Task Scheduling\nPassport.js Authentication", COLOR_PRIMARY, Inches(6.8), Inches(1.8)),
      ("Database & Security", "PostgreSQL 16 relational DB\nNeon Cloud Hosting\nPrisma ORM database driver\nJWT Access & Refresh token rotation\nGoogle OAuth 2.0 federation", COLOR_ACCENT, Inches(0.6), Inches(4.3)),
      ("Third-Party Integrations", "Razorpay Checkout Gateway\nResend Email Delivery API\nCloudinary Media CDN\nVercel (Frontend Delivery)\nRailway (Backend Container Hosting)", RGBColor(244, 63, 94), Inches(6.8), Inches(4.3))
    ]

    for title, desc, col, x, y in tech_data:
        # Card Background
        bg_shape = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.9), Inches(2.2))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = COLOR_CARD
        bg_shape.line.color.rgb = col
        bg_shape.line.width = Pt(1)
        
        # Textbox
        tb = slide3.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), Inches(5.5), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True
        
        tp = tf.paragraphs[0]
        tp.text = title
        tp.font.name = "Outfit"
        tp.font.size = Pt(16)
        tp.font.bold = True
        tp.font.color.rgb = col
        
        dp = tf.add_paragraph()
        dp.text = desc
        dp.font.name = "Inter"
        dp.font.size = Pt(10.5)
        dp.font.color.rgb = TEXT_WHITE
        dp.line_spacing = 1.25

    # ------------------------------------------
    # SLIDE 4: System Architecture
    # ------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_background(slide4)
    add_top_bar(slide4)
    add_header(slide4, "System Architecture Overview")
    add_footer(slide4, 4)

    # SVG Diagram Description
    descBox = slide4.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(4.8), Inches(4.8))
    tf = descBox.text_frame
    tf.word_wrap = True
    
    hp = tf.paragraphs[0]
    hp.text = "Decoupled Architecture Flow"
    hp.font.name = "Outfit"
    hp.font.size = Pt(20)
    hp.font.bold = True
    hp.font.color.rgb = COLOR_SECONDARY
    
    pts = [
      ("Modular Boundaries", "The backend is segmented into isolated domains with NestJS modules (Auth, Orders, Payments, Users)."),
      ("Stateless Guarding", "Next.js routes verify authenticated JWT states. Route access is controlled by custom NestJS role-based guards."),
      ("Cloud-Scale Storage", "Type-safe database interactions are handled via Prisma Client, connecting directly to PostgreSQL on Neon Cloud.")
    ]
    for pt_t, pt_d in pts:
        p = tf.add_paragraph()
        p.text = f"\n• {pt_t}: "
        p.font.name = "Inter"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
        dp = tf.add_paragraph()
        dp.text = pt_d
        dp.font.name = "Inter"
        dp.font.size = Pt(11)
        dp.font.color.rgb = TEXT_MUTED

    # Diagram drawing using PPTX shapes (representing Client, Backend, Storage, Ext Services)
    def draw_node(x, y, w, h, text, subtext, col):
        node = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        node.fill.solid()
        node.fill.fore_color.rgb = COLOR_CARD
        node.line.color.rgb = col
        node.line.width = Pt(1.5)
        
        tb = slide4.shapes.add_textbox(Inches(x), Inches(y + 0.1), Inches(w), Inches(h - 0.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Outfit"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        p.alignment = PP_ALIGN.CENTER
        
        sp = tf.add_paragraph()
        sp.text = subtext
        sp.font.name = "Inter"
        sp.font.size = Pt(8.5)
        sp.font.color.rgb = TEXT_MUTED
        sp.alignment = PP_ALIGN.CENTER

    # Clients Node
    draw_node(6.2, 1.8, 2.8, 1.0, "Next.js Applications", "Customer / Vendor / Super-Admin", COLOR_PRIMARY)
    # Backend Node
    draw_node(6.2, 3.4, 2.8, 1.5, "NestJS API & WS", "Authentication, Order Processors, Socket.IO WebSockets Server", COLOR_SECONDARY)
    # Database Node
    draw_node(6.2, 5.5, 2.8, 1.0, "PostgreSQL Database", "Neon Cloud + Prisma ORM", COLOR_ACCENT)
    
    # External Services
    draw_node(9.8, 2.5, 2.8, 1.2, "External Cloud Services", "Razorpay Payment Gateway\nResend Transactional Email\nGoogle OAuth Credentials", RGBColor(16, 185, 129))

    # Connectors (simple lines)
    def draw_line(x1, y1, x2, y2, col):
        connector = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x1), Inches(y1), Inches(x2-x1), Inches(y2-y1))
        connector.fill.solid()
        connector.fill.fore_color.rgb = col
        connector.line.fill.background()

    draw_line(7.55, 2.8, 7.65, 3.4, COLOR_PRIMARY) # Client to Backend
    draw_line(7.55, 4.9, 7.65, 5.5, COLOR_SECONDARY) # Backend to DB
    draw_line(9.0, 3.9, 9.8, 4.0, RGBColor(16, 185, 129)) # Backend to External

    # ------------------------------------------
    # SLIDE 5: Core Module A - Customer Portal
    # ------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_background(slide5)
    add_top_bar(slide5)
    add_header(slide5, "Module A: The Customer Portal")
    add_footer(slide5, 5)

    # 2 grid columns
    tbL = slide5.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.8))
    tfL = tbL.text_frame
    tfL.word_wrap = True
    
    hpL = tfL.paragraphs[0]
    hpL.text = "User Discovery & Checkout"
    hpL.font.name = "Outfit"
    hpL.font.size = Pt(20)
    hpL.font.bold = True
    hpL.font.color.rgb = COLOR_SECONDARY
    
    featuresL = [
      ("Dynamic Navigation strips", "Horizontally scrollable categories (Biryanis, Arabian, Beverages) with A-Z sort filters and Veg-Only selectors."),
      ("Local State Shopping Cart", "Maintains active cart items locally, validating and applying coupon promotional logic securely against backend REST endpoints."),
      ("Secured checkout modal", "Integrates native Razorpay checkout overlay processing payment signatures and confirming transactions.")
    ]
    for ft, fd in featuresL:
        p = tfL.add_paragraph()
        p.text = f"\n• {ft}: "
        p.font.name = "Inter"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
        dp = tfL.add_paragraph()
        dp.text = fd
        dp.font.name = "Inter"
        dp.font.size = Pt(11)
        dp.font.color.rgb = TEXT_MUTED

    # Right Content Card (WS Order Tracking)
    rcard = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.5))
    rcard.fill.solid()
    rcard.fill.fore_color.rgb = COLOR_CARD
    rcard.line.color.rgb = COLOR_ACCENT
    rcard.line.width = Pt(1.5)
    
    tbR = slide5.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.2), Inches(4.1))
    tfR = tbR.text_frame
    tfR.word_wrap = True
    
    hpR = tfR.paragraphs[0]
    hpR.text = "WebSocket-Driven Live Tracking"
    hpR.font.name = "Outfit"
    hpR.font.size = Pt(20)
    hpR.font.bold = True
    hpR.font.color.rgb = COLOR_ACCENT
    
    dpR1 = tfR.add_paragraph()
    dpR1.text = "\nReal-Time Order Tracking replaces manual browser refreshes with Socket.IO channels. Active status updates are pushed instantly from the kitchen dashboard."
    dpR1.font.name = "Inter"
    dpR1.font.size = Pt(13)
    dpR1.font.color.rgb = TEXT_WHITE
    
    dpR2 = tfR.add_paragraph()
    dpR2.text = "\nIntegrated Order Tracking Timeline:"
    dpR2.font.name = "Inter"
    dpR2.font.size = Pt(12)
    dpR2.font.bold = True
    dpR2.font.color.rgb = COLOR_SECONDARY
    
    dpR3 = tfR.add_paragraph()
    dpR3.text = "PENDING  ➔  CONFIRMED  ➔  PREPARING  ➔  OUT_FOR_DELIVERY  ➔  DELIVERED"
    dpR3.font.name = "Outfit"
    dpR3.font.size = Pt(11.5)
    dpR3.font.bold = True
    dpR3.font.color.rgb = RGBColor(16, 185, 129)

    # ------------------------------------------
    # SLIDE 6: Core Module B - Restaurant Admin
    # ------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_background(slide6)
    add_top_bar(slide6)
    add_header(slide6, "Module B: Restaurant Admin Console")
    add_footer(slide6, 6)

    # Left content column
    tbL6 = slide6.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.8))
    tfL6 = tbL6.text_frame
    tfL6.word_wrap = True
    
    hpL6 = tfL6.paragraphs[0]
    hpL6.text = "Store Operations & Live Boards"
    hpL6.font.name = "Outfit"
    hpL6.font.size = Pt(20)
    hpL6.font.bold = True
    hpL6.font.color.rgb = COLOR_SECONDARY
    
    bullets6 = [
      ("Live Incoming Order Board", "Real-time dispatch boards listen to incoming WebSockets, triggering immediate audio-visual alerts on screen."),
      ("Dynamic Menu Builder", "Perform full CRUD on dishes, adjust category tags, and toggle availability state dynamically."),
      ("Coupon Creator Panel", "Enables store admins to create discounts limited to their specific restaurant items.")
    ]
    for b_t, b_d in bullets6:
        p = tfL6.add_paragraph()
        p.text = f"\n• {b_t}: "
        p.font.name = "Inter"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
        dp = tfL6.add_paragraph()
        dp.text = b_d
        dp.font.name = "Inter"
        dp.font.size = Pt(11)
        dp.font.color.rgb = TEXT_MUTED

    # Right Content Column (Analytics Card)
    rcard6 = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.5))
    rcard6.fill.solid()
    rcard6.fill.fore_color.rgb = COLOR_CARD
    rcard6.line.color.rgb = COLOR_PRIMARY
    rcard6.line.width = Pt(1.5)
    
    tbR6 = slide6.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.2), Inches(4.1))
    tfR6 = tbR6.text_frame
    tfR6.word_wrap = True
    
    hpR6 = tfR6.paragraphs[0]
    hpR6.text = "Business Intelligence Dashboard"
    hpR6.font.name = "Outfit"
    hpR6.font.size = Pt(20)
    hpR6.font.bold = True
    hpR6.font.color.rgb = COLOR_PRIMARY
    
    dpR6 = tfR6.add_paragraph()
    dpR6.text = "\nIntegrated with Recharts, vendors monitor real-time sales and order metrics dynamically, including:"
    dpR6.font.name = "Inter"
    dpR6.font.size = Pt(12)
    dpR6.font.color.rgb = TEXT_WHITE
    
    # Sub stats
    stats = [
      ("Gross Analytics Sales", "₹ 24K+ (Overall platform orders)"),
      ("Orders Processed", "37 Completed Deliveries"),
      ("Revenue Trends", "Interactive 7-day revenue charts and popular categories donut charts.")
    ]
    for s_t, s_d in stats:
        p = tfR6.add_paragraph()
        p.text = f"• {s_t}: {s_d}"
        p.font.name = "Inter"
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_MUTED
        p.line_spacing = 1.3

    # ------------------------------------------
    # SLIDE 7: Core Module C - Super Admin
    # ------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    set_background(slide7)
    add_top_bar(slide7)
    add_header(slide7, "Module C: Super Admin Command Center")
    add_footer(slide7, 7)

    # 3 Column cards
    cards_data = [
      ("Platform Governance", "Central command dashboard provides total control over system status. Super administrators freeze, suspend, or reactivate customer or merchant accounts globally.", COLOR_SECONDARY, Inches(0.6)),
      ("Onboarding Pipelines", "All restaurant signups flow through a structured validation pipeline. Super admins audit application data and trigger status updates with simple interface buttons.", COLOR_PRIMARY, Inches(4.75)),
      ("System Audit Logging", "Every key administrative action—adjusting commission policy, modifying tax rates, blocking users—is written to database system logs for tracking and compliance.", COLOR_ACCENT, Inches(8.9))
    ]

    for title, desc, col, x in cards_data:
        card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(3.83), Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = col
        card.line.width = Pt(1.5)
        
        tb = slide7.shapes.add_textbox(x + Inches(0.15), Inches(2.0), Inches(3.53), Inches(4.1))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Outfit"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = col
        
        dp = tf.add_paragraph()
        dp.text = f"\n{desc}"
        dp.font.name = "Inter"
        dp.font.size = Pt(11.5)
        dp.font.color.rgb = TEXT_WHITE
        dp.line_spacing = 1.3

    # ------------------------------------------
    # SLIDE 8: COD Razorpay Bypass
    # ------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    set_background(slide8)
    add_top_bar(slide8)
    add_header(slide8, "Cash on Delivery (COD) Razorpay Bypass")
    add_footer(slide8, 8)

    # Left Column: Challenge & Resolution
    tbL8 = slide8.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.8))
    tfL8 = tbL8.text_frame
    tfL8.word_wrap = True
    
    hpL8 = tfL8.paragraphs[0]
    hpL8.text = "The Engineering Challenge"
    hpL8.font.name = "Outfit"
    hpL8.font.size = Pt(20)
    hpL8.font.bold = True
    hpL8.font.color.rgb = COLOR_SECONDARY
    
    bullets8 = [
      ("Razorpay Enforcement Bypass", "Bypasses the digital payment modal cleanly when a customer chooses COD, avoiding user-experience friction."),
      ("Conditional Routing (Next.js)", "Frontend intercepts COD selections, skips launching Razorpay SDK scripts, clears the shopping cart, and redirects instantly to order tracking."),
      ("Zero-Record Design Pattern", "Maintains clean relational tables. By skipping payment row instantiation for COD orders, database size is optimized and migrations are avoided.")
    ]
    for b_t, b_d in bullets8:
        p = tfL8.add_paragraph()
        p.text = f"\n• {b_t}: "
        p.font.name = "Inter"
        p.font.size = Pt(12.5)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
        dp = tfL8.add_paragraph()
        dp.text = b_d
        dp.font.name = "Inter"
        dp.font.size = Pt(11)
        dp.font.color.rgb = TEXT_MUTED

    # Right Column: Flow Chart Boxes
    steps = [
      ("1. Selection", "Customer selects Cash on Delivery on checkout."),
      ("2. Intercept", "Next.js bypasses Razorpay checkout SDK script."),
      ("3. Backend Placement", "NestJS stores order state directly as PENDING.")
    ]
    y_start = 1.8
    for num_t, desc_t in steps:
        step_bg = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(y_start), Inches(5.8), Inches(1.2))
        step_bg.fill.solid()
        step_bg.fill.fore_color.rgb = COLOR_CARD
        step_bg.line.color.rgb = COLOR_PRIMARY
        step_bg.line.width = Pt(1)
        
        tb = slide8.shapes.add_textbox(Inches(6.9), Inches(y_start + 0.1), Inches(5.6), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = num_t
        p.font.name = "Outfit"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_SECONDARY
        
        dp = tf.add_paragraph()
        dp.text = desc_t
        dp.font.name = "Inter"
        dp.font.size = Pt(11)
        dp.font.color.rgb = TEXT_WHITE
        
        y_start += 1.5

    # ------------------------------------------
    # SLIDE 9: Dynamic Invoicing
    # ------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    set_background(slide9)
    add_top_bar(slide9)
    add_header(slide9, "Dynamic Invoicing & Transactional Loop")
    add_footer(slide9, 9)

    # Left Column: Steps
    steps9 = [
      ("A. Order Delivered Trigger", "Store admin advances COD order to DELIVERED status."),
      ("B. Automated Capture", "NestJS updates order paymentStatus to PAID immediately."),
      ("C. Dynamic PDF Build", "InvoiceService builds receipt PDF labeling payment method as 'COD'."),
      ("D. Email Dispatch", "Resend API sends invoice PDF attachment to customer inbox.")
    ]
    y_s = 1.8
    for step_num, step_desc in steps9:
        s_bg = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(y_s), Inches(5.8), Inches(1.0))
        s_bg.fill.solid()
        s_bg.fill.fore_color.rgb = COLOR_CARD
        s_bg.line.color.rgb = COLOR_ACCENT
        s_bg.line.width = Pt(1)
        
        tb = slide9.shapes.add_textbox(Inches(0.7), Inches(y_s + 0.05), Inches(5.6), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = step_num
        p.font.name = "Outfit"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT
        
        dp = tf.add_paragraph()
        dp.text = step_desc
        dp.font.name = "Inter"
        dp.font.size = Pt(10.5)
        dp.font.color.rgb = TEXT_WHITE
        
        y_s += 1.25

    # Right Column: Closing the Loop description
    tbR9 = slide9.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.8))
    tfR9 = tbR9.text_frame
    tfR9.word_wrap = True
    
    hpR9 = tfR9.paragraphs[0]
    hpR9.text = "Closing the Operations Loop"
    hpR9.font.name = "Outfit"
    hpR9.font.size = Pt(20)
    hpR9.font.bold = True
    hpR9.font.color.rgb = COLOR_PRIMARY
    
    dpR9_1 = tfR9.add_paragraph()
    dpR9_1.text = "\nBy integrating invoice generation automatically, COD orders achieve parity with online credit card purchases. The system ensures transaction transparency without vendor manual intervention."
    dpR9_1.font.name = "Inter"
    dpR9_1.font.size = Pt(13)
    dpR9_1.font.color.rgb = TEXT_WHITE
    dpR9_1.line_spacing = 1.3
    
    # Mailbox mockup shape
    m_box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(4.2), Inches(5.8), Inches(2.1))
    m_box.fill.solid()
    m_box.fill.fore_color.rgb = RGBColor(27, 38, 59)
    m_box.line.color.rgb = RGBColor(16, 185, 129)
    m_box.line.width = Pt(1)
    
    tb_m = slide9.shapes.add_textbox(Inches(7.0), Inches(4.3), Inches(5.4), Inches(1.9))
    tf_m = tb_m.text_frame
    tf_m.word_wrap = True
    pm = tf_m.paragraphs[0]
    pm.text = "Email Delivery Mockup (Resend API)"
    pm.font.name = "Outfit"
    pm.font.size = Pt(13)
    pm.font.bold = True
    pm.font.color.rgb = RGBColor(16, 185, 129)
    
    pm_d = tf_m.add_paragraph()
    pm_d.text = "\nSubject: Your FoodFlow Invoice - Order Delivered\nAttachment: invoice_order_3728.pdf\nMessage: Thank you for ordering from FoodFlow! Your receipt is attached."
    pm_d.font.name = "Inter"
    pm_d.font.size = Pt(10.5)
    pm_d.font.color.rgb = TEXT_WHITE

    # ------------------------------------------
    # SLIDE 10: Technical Challenges
    # ------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    set_background(slide10)
    add_top_bar(slide10)
    add_header(slide10, "Technical Challenges & Resolutions")
    add_footer(slide10, 10)

    # 3 cards showing Challenge ➔ Resolution
    challenges = [
      ("Timezone Analytics", "Mismatched day boundaries caused empty gaps and date inaccuracies in dashboard charts.", "Standardized database operations to UTC and implemented dynamic timezone rendering client-side based on browser locations.", COLOR_SECONDARY, Inches(0.6)),
      ("Tenant Data Isolation", "Preventing data access crossovers and protecting sensitive sales statistics between store vendors.", "Developed custom NestJS guards to intercept incoming routes, verifying authenticated JWT token claims against query tenant targets.", COLOR_PRIMARY, Inches(4.75)),
      ("Node Environment Path", "Broken global symlinks on host profiles created server start and build path blockages.", "Configured explicit runtime variables in shell profiles to load localized appdata executable runtimes directly.", COLOR_ACCENT, Inches(8.9))
    ]

    for title, prob, sol, col, x in challenges:
        card = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(3.83), Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = col
        card.line.width = Pt(1.5)
        
        tb = slide10.shapes.add_textbox(x + Inches(0.15), Inches(1.9), Inches(3.53), Inches(4.3))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Outfit"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = col
        
        p_prob = tf.add_paragraph()
        p_prob.text = f"\nChallenge:\n{prob}"
        p_prob.font.name = "Inter"
        p_prob.font.size = Pt(10.5)
        p_prob.font.color.rgb = RGBColor(248, 113, 113) # red
        
        p_sol = tf.add_paragraph()
        p_sol.text = f"\nResolution:\n{sol}"
        p_sol.font.name = "Inter"
        p_sol.font.size = Pt(10.5)
        p_sol.font.color.rgb = RGBColor(74, 222, 128) # green
        p_sol.line_spacing = 1.25

    # ------------------------------------------
    # SLIDE 11: Learning Reflection
    # ------------------------------------------
    slide11 = prs.slides.add_slide(blank_layout)
    set_background(slide11)
    add_top_bar(slide11)
    add_header(slide11, "Internship Learning Reflections")
    add_footer(slide11, 11)

    # 4 Quadrants
    quads = [
      ("Enterprise System Architecture", "Gained rich knowledge in designing modular APIs, structuring decoupled databases, and managing robust micro-services with NestJS modules.", COLOR_SECONDARY, Inches(0.6), Inches(1.8)),
      ("Real-Time Communications", "Acquired practical skills in configuring WebSockets to transmit live updates and managing multi-client subscription rooms efficiently.", COLOR_PRIMARY, Inches(6.8), Inches(1.8)),
      ("External SDK Integrations", "Learned authentication flows (OAuth 2.0, JWT rotative tokens) and third-party systems integrations (Razorpay SDK, Resend mailing API).", COLOR_ACCENT, Inches(0.6), Inches(4.3)),
      ("Professional SDLC Pipeline", "Participated in practical developer operations: schema design, script seeding, environment settings isolation, and multi-cloud build pipelines.", RGBColor(16, 185, 129), Inches(6.8), Inches(4.3))
    ]

    for q_t, q_d, col, x, y in quads:
        bg_s = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.9), Inches(2.2))
        bg_s.fill.solid()
        bg_s.fill.fore_color.rgb = COLOR_CARD
        bg_s.line.color.rgb = col
        bg_s.line.width = Pt(1.5)
        
        tb = slide11.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), Inches(5.5), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = q_t
        p.font.name = "Outfit"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = col
        
        dp = tf.add_paragraph()
        dp.text = f"\n{q_d}"
        dp.font.name = "Inter"
        dp.font.size = Pt(11)
        dp.font.color.rgb = TEXT_WHITE
        dp.line_spacing = 1.3

    # ------------------------------------------
    # SLIDE 12: Future Scope & Conclusion
    # ------------------------------------------
    slide12 = prs.slides.add_slide(blank_layout)
    set_background(slide12)
    add_top_bar(slide12)
    add_header(slide12, "Future Scope & Project Conclusion")
    add_footer(slide12, 12)

    # Left Column: Future enhancements
    tbL12 = slide12.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.8))
    tfL12 = tbL12.text_frame
    tfL12.word_wrap = True
    
    hpL12 = tfL12.paragraphs[0]
    hpL12.text = "Future Scalability Paths"
    hpL12.font.name = "Outfit"
    hpL12.font.size = Pt(20)
    hpL12.font.bold = True
    hpL12.font.color.rgb = COLOR_SECONDARY
    
    enh = [
      ("AI Recommendations", "Incorporate personalized food item suggestions based on user ratings and ML predictive order models."),
      ("GPS Driver Dispatch", "Develop live mapping tracking using Mapbox to connect delivery agents and track courier paths on maps."),
      ("Subcription Plans", "Offer vendors subscription plans for premium dashboard layouts, shifting from base transaction commission fee margins.")
    ]
    for e_t, e_d in enh:
        p = tfL12.add_paragraph()
        p.text = f"\n• {e_t}: "
        p.font.name = "Inter"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
        dp = tfL12.add_paragraph()
        dp.text = e_d
        dp.font.name = "Inter"
        dp.font.size = Pt(11)
        dp.font.color.rgb = TEXT_MUTED

    # Right Column: Big final badge
    rcard12 = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.5))
    rcard12.fill.solid()
    rcard12.fill.fore_color.rgb = COLOR_CARD
    rcard12.line.color.rgb = COLOR_PRIMARY
    rcard12.line.width = Pt(1.5)
    
    tbR12 = slide12.shapes.add_textbox(Inches(7.1), Inches(2.2), Inches(5.2), Inches(3.8))
    tfR12 = tbR12.text_frame
    tfR12.word_wrap = True
    
    hpR12 = tfR12.paragraphs[0]
    hpR12.text = "FOODFLOW 4.1"
    hpR12.font.name = "Outfit"
    hpR12.font.size = Pt(28)
    hpR12.font.bold = True
    hpR12.font.color.rgb = COLOR_PRIMARY
    hpR12.alignment = PP_ALIGN.CENTER
    
    dpR12 = tfR12.add_paragraph()
    dpR12.text = "\n\nSuccessfully fulfills all B.Tech internship objectives, providing a scalable, secure, and production-ready SaaS food commerce platform."
    dpR12.font.name = "Inter"
    dpR12.font.size = Pt(14)
    dpR12.font.color.rgb = TEXT_WHITE
    dpR12.alignment = PP_ALIGN.CENTER
    dpR12.line_spacing = 1.3
    
    dpR12_2 = tfR12.add_paragraph()
    dpR12_2.text = "\n\n✔ Deployment & Validation Complete"
    dpR12_2.font.name = "Outfit"
    dpR12_2.font.size = Pt(13)
    dpR12_2.font.bold = True
    dpR12_2.font.color.rgb = RGBColor(16, 185, 129)
    dpR12_2.alignment = PP_ALIGN.CENTER

    prs.save("FOODFLOW_Presentation.pptx")
    print("PowerPoint presentation generated successfully as FOODFLOW_Presentation.pptx")

# ==========================================
# PDF GENERATOR USING REPORTLAB
# ==========================================
def create_pdfs():
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, KeepTogether
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors

    # --- STYLE SHEET BUILD ---
    styles = getSampleStyleSheet()
    
    # Custom Palette
    COLOR_PDF_NAVY = colors.HexColor('#0f172a')
    COLOR_PDF_CYAN = colors.HexColor('#0ea5e9')
    COLOR_PDF_TEXT = colors.HexColor('#1e293b')
    COLOR_PDF_LIGHT_BG = colors.HexColor('#f8fafc')
    COLOR_PDF_BORDER = colors.HexColor('#e2e8f0')

    title_style = ParagraphStyle(
        'DocTitle',
        fontName='Helvetica-Bold',
        fontSize=28,
        textColor=COLOR_PDF_NAVY,
        leading=34,
        alignment=0, # Left-aligned
        spaceAfter=15
    )

    subtitle_style = ParagraphStyle(
        'DocSubtitle',
        fontName='Helvetica',
        fontSize=14,
        textColor=colors.HexColor('#64748b'),
        leading=18,
        alignment=0,
        spaceAfter=30
    )

    h1_style = ParagraphStyle(
        'Heading1_Custom',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=18,
        textColor=COLOR_PDF_NAVY,
        leading=22,
        spaceBefore=18,
        spaceAfter=10,
        keepWithNext=True
    )

    h2_style = ParagraphStyle(
        'Heading2_Custom',
        parent=styles['Heading2'],
        fontName='Helvetica-Bold',
        fontSize=13,
        textColor=COLOR_PDF_CYAN,
        leading=16,
        spaceBefore=12,
        spaceAfter=6,
        keepWithNext=True
    )

    body_style = ParagraphStyle(
        'BodyText_Custom',
        parent=styles['BodyText'],
        fontName='Helvetica',
        fontSize=10,
        textColor=COLOR_PDF_TEXT,
        leading=14.5,
        spaceAfter=8
    )

    bullet_style = ParagraphStyle(
        'Bullet_Custom',
        parent=body_style,
        leftIndent=15,
        firstLineIndent=-10,
        spaceAfter=6
    )

    meta_label = ParagraphStyle(
        'MetaLabel',
        fontName='Helvetica-Bold',
        fontSize=9,
        textColor=colors.HexColor('#475569'),
        leading=11
    )

    meta_val = ParagraphStyle(
        'MetaVal',
        fontName='Helvetica',
        fontSize=9,
        textColor=COLOR_PDF_TEXT,
        leading=11
    )

    # ------------------------------------------
    # 1. GENERATE INTERNSHIP_REPORT.PDF
    # ------------------------------------------
    print("Generating INTERNSHIP_REPORT.pdf...")
    doc_ir = SimpleDocTemplate("INTERNSHIP_REPORT.pdf", pagesize=letter,
                              rightMargin=54, leftMargin=54, topMargin=54, bottomMargin=54)
    story_ir = []

    # Title Cover Elements
    story_ir.append(Spacer(1, 40))
    story_ir.append(Paragraph("INTERNSHIP PROJECT REPORT", ParagraphStyle('Badge', fontName='Helvetica-Bold', fontSize=10, textColor=COLOR_PDF_CYAN, spaceAfter=8)))
    story_ir.append(Paragraph("FOODFLOW 4.1", title_style))
    story_ir.append(Paragraph("Cloud-Based Multi-Vendor Food Ordering SaaS", subtitle_style))
    
    # Horiz separator line
    line_table = Table([[""]], colWidths=[504])
    line_table.setStyle(TableStyle([
        ('LINEBELOW', (0,0), (-1,-1), 2, COLOR_PDF_CYAN),
        ('BOTTOMPADDING', (0,0), (-1,-1), 0),
        ('TOPPADDING', (0,0), (-1,-1), 0),
    ]))
    story_ir.append(line_table)
    story_ir.append(Spacer(1, 20))

    # Abstract/Summary Intro
    summary_text = (
        "<b>Executive Summary:</b> This report details the architectural design, full-stack development, "
        "and cloud deployment of <b>FOODFLOW 4.1</b>, an enterprise-grade, multi-vendor food ordering "
        "Software-as-a-Service (SaaS) ecosystem. Bridging the gap between academic theory and practical "
        "software engineering, this project demonstrates end-to-end full-stack development. The platform connects "
        "customers, local restaurant vendors, and platform administrators through dedicated, isolated portals. "
        "It integrates secure token-based authentication, real-time WebSocket state synchronizations, Google OAuth "
        "identity federation, automated transaction emails, and the Razorpay payment gateway. This report outlines the "
        "system requirements, architectural patterns, technical implementation details, major features (including a "
        "specialized Cash on Delivery flow and dynamic invoicing), technical challenges overcome, and internship learning outcomes."
    )
    story_ir.append(Paragraph(summary_text, body_style))
    story_ir.append(Spacer(1, 150))

    # Meta Info block at cover bottom
    meta_data = [
        [Paragraph("Submitted By:", meta_label), Paragraph("Sreeharan M Anilkumar", meta_val),
         Paragraph("Organization:", meta_label), Paragraph("Bairuha Tech", meta_val)],
        [Paragraph("Register No:", meta_label), Paragraph("MEC23AIM006", meta_val),
         Paragraph("Mentors:", meta_label), Paragraph("Ambika Raj, Ressaan", meta_val)],
        [Paragraph("Program:", meta_label), Paragraph("B.Tech, Artificial Intelligence & ML", meta_val),
         Paragraph("Date:", meta_label), Paragraph("June 2026", meta_val)]
    ]
    t_meta = Table(meta_data, colWidths=[80, 172, 80, 172])
    t_meta.setStyle(TableStyle([
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ('BOTTOMPADDING', (0,0), (-1,-1), 6),
        ('TOPPADDING', (0,0), (-1,-1), 6),
        ('LINEBELOW', (0,-1), (-1,-1), 0.5, COLOR_PDF_BORDER),
    ]))
    story_ir.append(t_meta)
    
    story_ir.append(PageBreak())

    # --- PAGE 2: MAIN CONTENT ---
    story_ir.append(Paragraph("1. Introduction & Problem Domain", h1_style))
    story_ir.append(Paragraph("1.1 Problem Statement", h2_style))
    story_ir.append(Paragraph(
        "Traditional local food ordering systems suffer from fragmented user workflows, lack of scalability, "
        "and limited analytical capabilities. Small local restaurants and independent kitchens often lack the "
        "capital to invest in affordable, specialized software to manage online ordering, menu configuration, "
        "and sales auditing. Customers face disjointed order-tracking experiences, while platform owners struggle to "
        "govern vendor onboarding, enforce marketplace policies, or generate dynamic transactional documents.",
        body_style
    ))

    story_ir.append(Paragraph("1.2 Project Objectives", h2_style))
    story_ir.append(Paragraph("To address these bottlenecks, the objectives of this internship project were to:", body_style))
    story_ir.append(Paragraph("• <b>Build a Scalable Multi-Vendor Platform:</b> Design a decoupled, multi-tenant architecture allowing independent restaurant vendors to manage separate menus and orders.", bullet_style))
    story_ir.append(Paragraph("• <b>Provide Role-Based Portals:</b> Construct isolated user experiences for Customers, Restaurant Admins (Vendors), and Super Administrators.", bullet_style))
    story_ir.append(Paragraph("• <b>Enforce Security & Data Isolation:</b> Implement secure authentication (JWT with refresh rotation) and row-level multi-tenant database isolation.", bullet_style))
    story_ir.append(Paragraph("• <b>Enable Real-Time Operations:</b> Integrate WebSockets for live status updates and alerts on active orders.", bullet_style))
    story_ir.append(Paragraph("• <b>Integrate Transactional Systems:</b> Embed secure payments (Razorpay), automated email notifications (Resend), and dynamic invoicing.", bullet_style))

    story_ir.append(Spacer(1, 10))
    story_ir.append(Paragraph("2. Technology Stack & Architecture", h1_style))
    story_ir.append(Paragraph(
        "FOODFLOW 4.1 utilizes a decoupled client-server architecture. The system is designed to scale independently, "
        "maintaining separation of concerns between presentation, business logic, and storage layers.",
        body_style
    ))

    # Tech Stack Table
    table_data = [
        [Paragraph("<b>Layer</b>", meta_label), Paragraph("<b>Technologies</b>", meta_label), Paragraph("<b>Role / Rationale</b>", meta_label)],
        [Paragraph("Frontend", body_style), Paragraph("Next.js 16/15 (App Router), React 19, TypeScript, Tailwind CSS v4, Framer Motion", body_style), Paragraph("Deliver a fast, responsive, and SEO-optimized UI with smooth transitions.", body_style)],
        [Paragraph("Backend", body_style), Paragraph("NestJS 11, Node.js, Socket.IO, Passport.js", body_style), Paragraph("Provide a modular, structured gateway and WebSocket server.", body_style)],
        [Paragraph("Database", body_style), Paragraph("PostgreSQL 16 (Neon Cloud), Prisma ORM", body_style), Paragraph("Ensure relational storage, type safety, and transactional control.", body_style)],
        [Paragraph("Security", body_style), Paragraph("JWT (Token Rotation), Google OAuth 2.0", body_style), Paragraph("Enforce secure stateless auth and social identity federation.", body_style)],
        [Paragraph("Integrations", body_style), Paragraph("Razorpay SDK, Resend API, Cloudinary", body_style), Paragraph("Handle checkout checkout payments, emails, and asset CDNs.", body_style)]
    ]
    t_stack = Table(table_data, colWidths=[80, 180, 244])
    t_stack.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), COLOR_PDF_LIGHT_BG),
        ('GRID', (0,0), (-1,-1), 0.5, COLOR_PDF_BORDER),
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ('PADDING', (0,0), (-1,-1), 6),
    ]))
    story_ir.append(t_stack)

    story_ir.append(PageBreak())

    # --- PAGE 3: DELIVERABLES ---
    story_ir.append(Paragraph("3. System Implementation & Deliverables", h1_style))
    story_ir.append(Paragraph("3.1 Database Design", h2_style))
    story_ir.append(Paragraph(
        "The database structure leverages a highly normalized PostgreSQL schema. The core entities include: "
        "<b>User</b> (distinct roles: <i>CUSTOMER, ADMIN, SUPER_ADMIN</i>), <b>Restaurant</b> (onboarding states), "
        "<b>Food & Category</b> (hierarchical menus), <b>Order & OrderItem</b> (purchases snapshot), "
        "<b>Payment & Invoice</b> (financial audits), and <b>AuditLog</b> (system-wide administrative tracking).",
        body_style
    ))

    story_ir.append(Paragraph("3.2 System Modules", h2_style))
    story_ir.append(Paragraph(
        "<b>Module A: The Customer Portal</b> - Tailored to discovery and checkout. Includes horizontally scrollable categories, "
        "smart cart coupon validation, and WebSocket-driven live order tracking timeline updates.",
        bullet_style
    ))
    story_ir.append(Paragraph(
        "<b>Module B: The Restaurant Admin Console</b> - A vendor dashboard with multi-tenant isolation. Includes live incoming "
        "order queues, menu builders, and Recharts-based 7-day revenue trend visualizations.",
        bullet_style
    ))
    story_ir.append(Paragraph(
        "<b>Module C: The Super Admin Control Center</b> - The central command console. Tracks gross platform sales, active merchants, "
        "registered user records, and logs system-wide administrative audit actions.",
        bullet_style
    ))

    story_ir.append(Spacer(1, 10))
    story_ir.append(Paragraph("4. Key Engineering Deliverables", h1_style))
    story_ir.append(Paragraph("4.1 Cash on Delivery (COD) Razorpay Bypass", h2_style))
    story_ir.append(Paragraph(
        "To accommodate door-step cash payments without database migrations, we implemented a bypass: "
        "1. <b>Frontend Integration:</b> Next.js intercepts COD selections, bypasses Razorpay SDK script overlay execution, and routes directly to the tracking screen. "
        "2. <b>Backend Logic:</b> NestJS skips payment initialization, storing the order state as PENDING immediately. "
        "3. <b>Data Integrity:</b> A zero-record design pattern was utilized, skipping payment row instantiation for COD orders to preserve database structure.",
        body_style
    ))

    story_ir.append(Paragraph("4.2 Dynamic Invoicing & Lifecycle Completion", h2_style))
    story_ir.append(Paragraph(
        "COD orders complete their lifecycle as follows: "
        "1. <b>Status Capture:</b> When the store admin updates the order status to DELIVERED, the backend flags paymentStatus to PAID. "
        "2. <b>PDF Build:</b> InvoiceService dynamically generates a receipt showing Cash on Delivery as the method, and N/A for transaction ID. "
        "3. <b>Mailing:</b> The receipt is sent immediately as an email attachment to the customer via the Resend API.",
        body_style
    ))

    story_ir.append(PageBreak())

    # --- PAGE 4: CHALLENGES & REFLECTION ---
    story_ir.append(Paragraph("5. Technical Challenges & Solutions", h1_style))
    story_ir.append(Paragraph("• <b>Timezone-Safe Analytics:</b> Solved chart boundary discrepancies by standardizing database timestamps to UTC and formatting client-side graphs relative to local user location.", bullet_style))
    story_ir.append(Paragraph("• <b>Tenant Data Isolation:</b> Created NestJS Guards that cross-verify authenticated JWT claims against request path parameters, preventing vendors from accessing another merchant's financial logs.", bullet_style))
    story_ir.append(Paragraph("• <b>Node Environment PATH:</b> Resolved path executable conflicts on the host by defining explicit system shell runtime path profile variables.", bullet_style))

    story_ir.append(Spacer(1, 15))
    story_ir.append(Paragraph("6. Internship Learning Outcomes & Conclusion", h1_style))
    story_ir.append(Paragraph("6.1 Practical Skills Acquired", h2_style))
    story_ir.append(Paragraph("• <b>Enterprise Architecture:</b> Mastered NestJS modular structure, dependency injection, and Prisma client transactions.", bullet_style))
    story_ir.append(Paragraph("• <b>Real-Time System Design:</b> Handled WebSocket channel room management, Socket.IO clients, and cache invalidation.", bullet_style))
    story_ir.append(Paragraph("• <b>Integrations & Security:</b> Integrated Razorpay SDK checkout scripts, Google OAuth federation, and Resend mail APIs.", bullet_style))
    story_ir.append(Paragraph("• <b>Operations & Delivery:</b> Conducted database schema migrations, environment isolation, and Railway/Vercel continuous integration pipelines.", bullet_style))

    story_ir.append(Paragraph("6.2 Conclusion", h2_style))
    story_ir.append(Paragraph(
        "The Bairuha Tech internship successfully delivered FOODFLOW 4.1 as a robust, scalable, multi-vendor SaaS platform. "
        "Taking the codebase from concept to fully synchronized cloud deployment deepened my understanding of web engineering, "
        "data security, and transactional reliability. I am confident that these engineering insights will serve as a strong "
        "foundation for my future career in full-stack cloud systems engineering.",
        body_style
    ))

    # Build PDF
    doc_ir.build(story_ir)
    print("INTERNSHIP_REPORT.pdf built successfully!")

    # ------------------------------------------
    # 2. GENERATE SELF_LEARNING_REPORT.PDF (Must fit exactly on 1 page!)
    # ------------------------------------------
    print("Generating SELF_LEARNING_REPORT.pdf (1-page restriction)...")
    doc_sl = SimpleDocTemplate("SELF_LEARNING_REPORT.pdf", pagesize=letter,
                              rightMargin=45, leftMargin=45, topMargin=45, bottomMargin=45)
    story_sl = []

    # Title & Student details
    story_sl.append(Paragraph("Internship Learning Reflection Report", ParagraphStyle('RefTitle', fontName='Helvetica-Bold', fontSize=22, textColor=COLOR_PDF_NAVY, leading=26, spaceAfter=15)))
    
    # Meta Block
    meta_sl = [
        [Paragraph("<b>Student Name:</b>", meta_label), Paragraph("Sreeharan M Anilkumar", meta_val),
         Paragraph("<b>Register No:</b>", meta_label), Paragraph("MEC23AIM006", meta_val)],
        [Paragraph("<b>Program:</b>", meta_label), Paragraph("B.Tech, Artificial Intelligence & ML", meta_val),
         Paragraph("<b>Organization:</b>", meta_label), Paragraph("Bairuha Tech", meta_val)],
        [Paragraph("<b>Mentors:</b>", meta_label), Paragraph("Ambika Raj, Ressaan", meta_val),
         Paragraph("<b>Project:</b>", meta_label), Paragraph("FOODFLOW 4.1 – SaaS Food Ordering", meta_val)]
    ]
    t_sl = Table(meta_sl, colWidths=[90, 172, 80, 180])
    t_sl.setStyle(TableStyle([
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('BOTTOMPADDING', (0,0), (-1,-1), 4),
        ('TOPPADDING', (0,0), (-1,-1), 4),
        ('LINEBELOW', (0,-1), (-1,-1), 1, COLOR_PDF_CYAN),
    ]))
    story_sl.append(t_sl)
    story_sl.append(Spacer(1, 15))

    # Reflection Body Sections
    story_sl.append(Paragraph("1. Introduction and Objectives", h2_style))
    story_sl.append(Paragraph(
        "During my internship at <b>Bairuha Tech</b>, I designed, developed, and deployed FOODFLOW 4.1, a "
        "cloud-based multi-vendor food ordering platform connecting customers, restaurant vendors, and administrators. "
        "The project bridged academic B.Tech AI & ML knowledge with real-world software engineering practices, providing "
        "exposure to enterprise-scale application development, API design, and multi-tenant security structures.",
        body_style
    ))

    story_sl.append(Paragraph("2. Technical Skills Acquired", h2_style))
    story_sl.append(Paragraph(
        "I gained hands-on full-stack development experience utilizing: <b>Next.js 16/15</b> (React 19, TypeScript), "
        "<b>Tailwind CSS v4</b>, and <b>Framer Motion</b> on the frontend; <b>NestJS 11</b>, <b>Node.js</b>, and "
        "<b>Socket.IO</b> for backend WebSockets APIs; and <b>PostgreSQL on Neon Cloud</b> with <b>Prisma ORM</b> for databases. "
        "I integrated <b>Razorpay</b> (payments), <b>Resend API</b> (emails), <b>Google OAuth 2.0</b>, and managed deployments using <b>Vercel</b> and <b>Railway</b>.",
        body_style
    ))

    story_sl.append(Paragraph("3. UI/UX Design and Architecture", h2_style))
    story_sl.append(Paragraph(
        "We followed a wireframe-first approach to validate user flows before writing code. I designed separate, isolated portal "
        "experiences for: <b>Customers</b> (menu discovery, live checkout), <b>Restaurant Admins</b> (live order board, CRUD, coupons), "
        "and <b>Super Administrators</b> (merchant approvals, lock accounts, audit logs) to optimize system operations.",
        body_style
    ))

    story_sl.append(Paragraph("4. Key Challenges and Problem Solving", h2_style))
    story_sl.append(Paragraph(
        "Core challenges resolved during this internship included: <b>1) Cash on Delivery flow integration:</b> created a customized checkout "
        "bypass skipping the Razorpay transaction overlay while preserving table schemas. <b>2) Tenant Data Isolation:</b> structured NestJS Guards "
        "matching token claims with route params to prevent cross-vendor metrics crossover. <b>3) Deployment configurations:</b> handled "
        "broken environment execution pathing directly on the host console.",
        body_style
    ))

    story_sl.append(Paragraph("5. Overall Reflection", h2_style))
    story_sl.append(Paragraph(
        "This internship provided end-to-end exposure to the software lifecycle—from database normalization and schema design to "
        "API implementation, testing, and multi-cloud deployment pipelines. Building FOODFLOW 4.1 has significantly strengthened "
        "my software engineering capability, preparing me for complex challenges in full-stack cloud system development.",
        body_style
    ))

    # Build PDF (fits strictly on 1 page!)
    doc_sl.build(story_sl)
    print("SELF_LEARNING_REPORT.pdf built successfully!")

# ==========================================
# MAIN EXECUTION ROUTINE
# ==========================================
if __name__ == "__main__":
    create_presentation()
    try:
        create_pdfs()
    except Exception as e:
        print(f"Error while generating PDFs: {e}", file=sys.stderr)
        print("Make sure reportlab is installed: pip install reportlab", file=sys.stderr)
